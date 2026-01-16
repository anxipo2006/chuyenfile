import streamlit as st
import os
import tempfile
import platform
import subprocess
import sys
from PIL import Image
from pdf2docx import Converter
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from deep_translator import GoogleTranslator

# --- CẤU HÌNH ---
st.set_page_config(page_title="Universal Tool: Convert & Translate", page_icon="☁️", layout="wide")

# --- CSS ---
st.markdown("""
<style>
    .main { background-color: #0e1117; font-family: 'Segoe UI', sans-serif; }
    .css-card { background-color: #262730; border-radius: 12px; padding: 20px; border: 1px solid #444; text-align: center; margin-bottom: 20px; }
    .stButton>button { width: 100%; font-weight: bold; border-radius: 8px; }
    .card-header { color: #00e676; font-weight: bold; font-size: 1.1em; margin-bottom: 10px; }
</style>
""", unsafe_allow_html=True)

# --- HÀM HỆ THỐNG ---
def save_uploaded_file(uploaded_file):
    try:
        suffix = f".{uploaded_file.name.split('.')[-1]}"
        with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
            tmp.write(uploaded_file.getbuffer())
            return tmp.name
    except: return None

def get_platform():
    return platform.system()

# --- HÀM CONVERT (HỖ TRỢ CẢ LINUX VÀ WINDOWS) ---

def convert_with_libreoffice(input_path, output_format="pdf"):
    """Chạy lệnh LibreOffice trên Linux để convert"""
    try:
        out_folder = os.path.dirname(input_path)
        # Lệnh gọi LibreOffice chế độ không giao diện (headless)
        cmd = [
            "libreoffice", "--headless", "--convert-to", output_format,
            input_path, "--outdir", out_folder
        ]
        subprocess.run(cmd, check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
        
        # Tên file đầu ra dự kiến
        filename = os.path.basename(input_path)
        name_no_ext = os.path.splitext(filename)[0]
        return os.path.join(out_folder, f"{name_no_ext}.{output_format}")
    except Exception as e:
        print(f"Lỗi LibreOffice: {e}")
        return None

def process_office_to_pdf(path):
    """Hàm thông minh tự chọn cách convert tùy theo hệ điều hành"""
    os_name = get_platform()
    
    # 1. Nếu là Linux (Trên Web Streamlit Cloud)
    if os_name == "Linux":
        return convert_with_libreoffice(path, "pdf")
    
    # 2. Nếu là Windows (Chạy máy cá nhân)
    elif os_name == "Windows":
        try:
            # Word
            if path.endswith(".docx") or path.endswith(".doc"):
                from docx2pdf import convert
                pdf_path = path.replace(".docx", ".pdf").replace(".doc", ".pdf")
                convert(path, pdf_path)
                return pdf_path
            
            # PPT
            elif path.endswith(".pptx") or path.endswith(".ppt"):
                import comtypes.client
                pdf_path = path.replace(".pptx", ".pdf").replace(".ppt", ".pdf")
                path = os.path.abspath(path)
                pdf_path = os.path.abspath(pdf_path)
                ppt = comtypes.client.CreateObject("Powerpoint.Application")
                ppt.Visible = 1
                deck = ppt.Presentations.Open(path)
                deck.SaveAs(pdf_path, 32)
                deck.Close()
                return pdf_path
        except Exception as e:
            # Nếu Windows mà không cài Office thì thử dùng LibreOffice (nếu có cài)
            # Hoặc báo lỗi
            return convert_with_libreoffice(path, "pdf")
            
    return None

def process_img_to_pdf(files):
    try:
        imgs = [Image.open(f).convert('RGB') for f in files]
        tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".pdf")
        tmp.close()
        imgs[0].save(tmp.name, save_all=True, append_images=imgs[1:])
        return tmp.name
    except: return None

def process_pdf_to_word(path):
    try:
        docx = path.replace(".pdf", ".docx")
        cv = Converter(path)
        cv.convert(docx)
        cv.close()
        return docx
    except: return None

# --- LOGIC DỊCH THUẬT (GIỮ NGUYÊN - VÌ NÓ CHẠY TỐT TRÊN MỌI OS) ---
translator = GoogleTranslator(source='auto', target='vi')

def safe_trans(text):
    if not text or not isinstance(text, str) or len(text) < 2 or text.isnumeric(): return text
    try: return translator.translate(text)
    except: return text

def process_trans_word(path, bar):
    doc = Document(path)
    total = len(doc.paragraphs)
    for i, p in enumerate(doc.paragraphs):
        if p.text.strip(): p.text = safe_trans(p.text)
        if i % 10 == 0: bar.progress(min(i/total, 0.9))
    
    for t in doc.tables:
        for r in t.rows:
            for c in r.cells:
                for p in c.paragraphs:
                     if p.text.strip(): p.text = safe_trans(p.text)
    
    bar.progress(1.0)
    out = path.replace(".docx", "_VN.docx")
    doc.save(out)
    return out

def process_trans_excel(path, bar):
    wb = load_workbook(path)
    sheets = wb.worksheets
    for i, s in enumerate(sheets):
        bar.progress(i/len(sheets))
        for r in s.iter_rows():
            for c in r:
                if c.value and isinstance(c.value, str) and not c.value.startswith("="):
                    c.value = safe_trans(c.value)
    bar.progress(1.0)
    out = path.replace(".xlsx", "_VN.xlsx")
    wb.save(out)
    return out

def process_trans_ppt(path, bar):
    prs = Presentation(path)
    slides = prs.slides
    
    def proc_shp(shp):
        if shp.shape_type == MSO_SHAPE_TYPE.GROUP:
            for child in shp.shapes: proc_shp(child)
        if shp.has_text_frame:
            for p in shp.text_frame.paragraphs:
                if p.text.strip():
                    orig = p.text
                    trans = safe_trans(orig)
                    p.text = trans 
        if shp.has_table:
            for r in shp.table.rows:
                for c in r.cells:
                    if c.text_frame:
                         for p in c.text_frame.paragraphs:
                             if p.text.strip(): p.text = safe_trans(p.text)

    for i, slide in enumerate(slides):
        bar.progress(i/len(slides))
        for shp in slide.shapes: proc_shp(shp)
        
    bar.progress(1.0)
    out = path.replace(".pptx", "_VN.pptx")
    prs.save(out)
    return out

# --- GIAO DIỆN ---
st.markdown("<h2 style='text-align:center;'>🌍 WEB DỊCH THUẬT & CHUYỂN ĐỔI (ONLINE)</h2>", unsafe_allow_html=True)

if get_platform() == "Linux":
    st.caption("✅ Đang chạy trên Server Linux (Streamlit Cloud). Sử dụng LibreOffice cho Converter.")
else:
    st.caption("✅ Đang chạy trên Windows Local.")

tab1, tab2 = st.tabs(["CHUYỂN ĐỔI FILE", "DỊCH THUẬT"])

with tab1:
    c1, c2 = st.columns(2)
    with c1:
        st.markdown("<div class='css-card'><div class='card-header'>Word/PPT ➡ PDF</div></div>", unsafe_allow_html=True)
        f = st.file_uploader("Chọn Word hoặc PPT", type=['docx', 'pptx'])
        if f and st.button("Chuyển sang PDF"):
            with st.spinner("Đang xử lý trên Cloud..."):
                p = save_uploaded_file(f)
                res = process_office_to_pdf(p)
                if res and os.path.exists(res):
                    with open(res, "rb") as file: st.download_button("Tải PDF", file, "converted.pdf")
                else:
                    st.error("Lỗi chuyển đổi. Trên Web, định dạng phức tạp có thể bị lỗi.")

    with c2:
        st.markdown("<div class='css-card'><div class='card-header'>Image ➡ PDF</div></div>", unsafe_allow_html=True)
        imgs = st.file_uploader("Chọn Ảnh", type=['png','jpg'], accept_multiple_files=True)
        if imgs and st.button("Gộp ảnh thành PDF"):
            res = process_img_to_pdf(imgs)
            if res:
                with open(res, "rb") as file: st.download_button("Tải PDF", file, "images.pdf")
    
    st.markdown("---")
    st.markdown("<div class='css-card'><div class='card-header'>PDF ➡ Word</div></div>", unsafe_allow_html=True)
    pdf = st.file_uploader("Chọn PDF", type=['pdf'])
    if pdf and st.button("Chuyển sang Word"):
        p = save_uploaded_file(pdf)
        res = process_pdf_to_word(p)
        if res: with open(res, "rb") as f: st.download_button("Tải Word", f, "converted.docx")

with tab2:
    st.info("Dịch thuật giữ nguyên format (Word, Excel, PPT)")
    col_t1, col_t2, col_t3 = st.columns(3)
    
    with col_t1:
        w_t = st.file_uploader("Word", type=['docx'], key='tw')
        if w_t and st.button("Dịch Word"):
            p = save_uploaded_file(w_t)
            res = process_trans_word(p, st.empty())
            with open(res, "rb") as f: st.download_button("Download", f, "VN_doc.docx")
            
    with col_t2:
        e_t = st.file_uploader("Excel", type=['xlsx'], key='te')
        if e_t and st.button("Dịch Excel"):
            p = save_uploaded_file(e_t)
            res = process_trans_excel(p, st.empty())
            with open(res, "rb") as f: st.download_button("Download", f, "VN_excel.xlsx")
            
    with col_t3:
        p_t = st.file_uploader("PPT", type=['pptx'], key='tp')
        if p_t and st.button("Dịch PPT"):
            p = save_uploaded_file(p_t)
            res = process_trans_ppt(p, st.empty())
            with open(res, "rb") as f: st.download_button("Download", f, "VN_slide.pptx")